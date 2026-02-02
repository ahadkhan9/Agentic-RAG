"""
Document Loader Module (Enhanced)

Handles loading and extracting text from various document formats:
- PDF (PyMuPDF) - with table detection and header/footer removal
- DOCX (python-docx) - with section hierarchy
- Excel (openpyxl) - with table-aware chunking
- PowerPoint (python-pptx) - with slide content extraction
- Text files

Data Engineering Features:
- Table extraction preserved as markdown
- Header/footer noise removal
- Section hierarchy tracking
- Content type classification
"""
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation

from logger import get_logger

logger = get_logger("Loader")


@dataclass
class DocumentChunk:
    """Represents a chunk of document content with rich metadata."""
    content: str
    source_file: str
    file_type: str
    page_number: Optional[int] = None
    section: Optional[str] = None
    section_path: list[str] = field(default_factory=list)  # Full hierarchy
    row_index: Optional[int] = None
    slide_number: Optional[int] = None
    content_type: str = "text"  # text, table, list, heading


# Common header/footer patterns to remove
NOISE_PATTERNS = [
    r'Page \d+ of \d+',  # Page X of Y
    r'^\s*\d+\s*$',  # Standalone page numbers
    r'Confidential.*?(?=\n|$)',  # Confidentiality notices
    r'Copyright ©.*?(?=\n|$)',  # Copyright notices
    r'^\s*-\s*\d+\s*-\s*$',  # Page markers like "- 5 -"
    r'Document ID:?\s*[\w-]+',  # Document IDs
    r'Revision:?\s*[\d.]+',  # Revision numbers
    r'^\s*Draft\s*$',  # Draft watermarks
    r'CONFIDENTIAL',  # All caps confidential
    r'INTERNAL USE ONLY',  # Internal notices
]


def clean_text(text: str) -> str:
    """Remove headers, footers, and noise from extracted text."""
    for pattern in NOISE_PATTERNS:
        text = re.sub(pattern, '', text, flags=re.MULTILINE | re.IGNORECASE)
    
    # Remove excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {3,}', ' ', text)
    
    return text.strip()


def detect_table_in_text(text: str) -> bool:
    """Heuristic to detect if text looks like a table."""
    lines = text.strip().split('\n')
    if len(lines) < 2:
        return False
    
    # Check for consistent column separators (tabs or multiple spaces)
    tab_lines = sum(1 for line in lines if '\t' in line or '  ' in line)
    return tab_lines > len(lines) * 0.5


def text_to_markdown_table(text: str) -> str:
    """Attempt to convert tab/space-separated text to markdown table."""
    lines = text.strip().split('\n')
    if len(lines) < 2:
        return text
    
    # Split by tabs or multiple spaces
    rows = []
    for line in lines:
        if '\t' in line:
            cells = [c.strip() for c in line.split('\t')]
        else:
            cells = [c.strip() for c in re.split(r'  +', line)]
        if any(cells):
            rows.append(cells)
    
    if not rows:
        return text
    
    # Normalize column count
    max_cols = max(len(row) for row in rows)
    rows = [row + [''] * (max_cols - len(row)) for row in rows]
    
    # Build markdown
    md = "| " + " | ".join(rows[0]) + " |\n"
    md += "| " + " | ".join("---" for _ in rows[0]) + " |\n"
    for row in rows[1:]:
        md += "| " + " | ".join(row) + " |\n"
    
    return md


def load_pdf(file_path: str) -> list[DocumentChunk]:
    """
    Load and extract text from a PDF file.
    Enhanced with:
    - Table detection
    - Header/footer removal
    - Content type classification
    """
    chunks = []
    filename = os.path.basename(file_path)
    
    logger.info(f"📄 Loading PDF: {filename}")
    
    doc = fitz.open(file_path)
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        
        if not text:
            continue
        
        # Clean headers/footers
        text = clean_text(text)
        
        if not text:
            continue
        
        # Detect if this looks like a table
        if detect_table_in_text(text):
            logger.debug(f"  Page {page_num}: Detected table content")
            chunks.append(DocumentChunk(
                content=text_to_markdown_table(text),
                source_file=filename,
                file_type="pdf",
                page_number=page_num,
                content_type="table"
            ))
        else:
            chunks.append(DocumentChunk(
                content=text,
                source_file=filename,
                file_type="pdf",
                page_number=page_num,
                content_type="text"
            ))
    
    doc.close()
    logger.info(f"  ✅ Extracted {len(chunks)} chunks from {filename}")
    
    return chunks


def load_docx(file_path: str) -> list[DocumentChunk]:
    """
    Load and extract text from a DOCX file.
    Enhanced with:
    - Section hierarchy tracking
    - Heading detection
    - Content type classification
    """
    chunks = []
    filename = os.path.basename(file_path)
    
    logger.info(f"📝 Loading DOCX: {filename}")
    
    doc = DocxDocument(file_path)
    section_stack = []  # Full hierarchy: ["Chapter 1", "Safety", "LOTO"]
    current_section = "Document"
    page_estimate = 1
    char_count = 0
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        
        # Clean noise
        text = clean_text(text)
        if not text:
            continue
        
        # Track heading hierarchy
        style_name = para.style.name if para.style else ""
        content_type = "text"
        
        if style_name.startswith('Heading'):
            level = int(style_name.replace('Heading ', '') or 1)
            # Pop stack to current level
            section_stack = section_stack[:level - 1]
            section_stack.append(text)
            current_section = text
            content_type = "heading"
        
        # Estimate page (rough: ~3000 chars per page)
        char_count += len(text)
        page_estimate = (char_count // 3000) + 1
        
        chunks.append(DocumentChunk(
            content=text,
            source_file=filename,
            file_type="docx",
            page_number=page_estimate,
            section=current_section,
            section_path=section_stack.copy(),
            content_type=content_type
        ))
    
    # Also extract tables from DOCX
    for table_idx, table in enumerate(doc.tables):
        table_text = extract_docx_table(table)
        if table_text:
            chunks.append(DocumentChunk(
                content=table_text,
                source_file=filename,
                file_type="docx",
                section=f"Table {table_idx + 1}",
                content_type="table"
            ))
    
    logger.info(f"  ✅ Extracted {len(chunks)} chunks from {filename}")
    
    return chunks


def extract_docx_table(table) -> str:
    """Extract DOCX table as markdown."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows.append(cells)
    
    if not rows:
        return ""
    
    # Build markdown
    md = "| " + " | ".join(rows[0]) + " |\n"
    md += "| " + " | ".join("---" for _ in rows[0]) + " |\n"
    for row in rows[1:]:
        md += "| " + " | ".join(row) + " |\n"
    
    return md


def load_excel(file_path: str) -> list[DocumentChunk]:
    """
    Load and extract data from an Excel file.
    Enhanced with:
    - Auto header row detection
    - Semantic text conversion (natural language per row)
    - Multiple chunking strategies
    - File context preservation for search
    """
    chunks = []
    filename = os.path.basename(file_path)
    
    logger.info(f"📊 Loading Excel: {filename}")
    
    excel_file = pd.ExcelFile(file_path)
    
    for sheet_name in excel_file.sheet_names:
        # Try to detect header row (often row 0 has title, row 1 has headers)
        # Read first few rows to detect proper header
        df_preview = pd.read_excel(excel_file, sheet_name=sheet_name, nrows=5, header=None)
        
        # Find the header row by looking for row with most string values
        header_row = 0
        for idx in range(min(3, len(df_preview))):
            row = df_preview.iloc[idx]
            # Check if this row looks like headers (mostly strings, no pure numbers as column names)
            if row.notna().sum() >= 3:
                # Check if values look like column headers
                has_unnamed = any('unnamed' in str(v).lower() for v in row if pd.notna(v))
                if not has_unnamed:
                    header_row = idx
                    break
        
        # Read with detected header
        df = pd.read_excel(excel_file, sheet_name=sheet_name, header=header_row)
        
        if df.empty:
            continue
        
        # Clean column names
        df.columns = [str(c).strip() for c in df.columns]
        
        # Strategy 1: Full table as markdown (always create this for overview)
        table_content = f"**Excel File: {filename}**\n"
        table_content += f"**Sheet: {sheet_name}**\n\n"
        table_content += df.to_markdown(index=False)
        
        chunks.append(DocumentChunk(
            content=table_content,
            source_file=filename,
            file_type="excel",
            section=f"{sheet_name} (Full Table)",
            content_type="table"
        ))
        
        # Strategy 2: Create semantic text per row for better retrieval
        headers = df.columns.tolist()
        
        for idx, row in df.iterrows():
            # Skip rows with mostly NaN values
            valid_values = sum(1 for v in row if pd.notna(v))
            if valid_values < 2:
                continue
            
            # Build natural language description
            row_text = f"From Excel file '{filename}', sheet '{sheet_name}':\n"
            
            # Add key-value pairs as natural text
            for col in headers:
                value = row[col]
                if pd.notna(value):
                    row_text += f"- {col}: {value}\n"
            
            if row_text.strip() and len(row_text) > 50:  # Only meaningful content
                chunks.append(DocumentChunk(
                    content=row_text.strip(),
                    source_file=filename,
                    file_type="excel",
                    row_index=int(idx) + header_row + 2,  # Excel row number
                    section=sheet_name,
                    content_type="table_row"
                ))
    
    logger.info(f"  ✅ Extracted {len(chunks)} chunks from {filename}")
    
    return chunks


def load_pptx(file_path: str) -> list[DocumentChunk]:
    """
    Load and extract text from a PowerPoint file.
    Enhanced with:
    - Speaker notes extraction
    - Slide title tracking
    - Content type classification
    """
    chunks = []
    filename = os.path.basename(file_path)
    
    logger.info(f"📽️ Loading PowerPoint: {filename}")
    
    prs = Presentation(file_path)
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        slide_title = f"Slide {slide_num}"
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = clean_text(shape.text.strip())
                if text:
                    slide_text.append(text)
                    # First text is usually the title
                    if len(slide_text) == 1:
                        slide_title = text[:50]
        
        # Extract speaker notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            notes = clean_text(notes)
            if notes:
                slide_text.append(f"[Speaker Notes: {notes}]")
        
        if slide_text:
            chunks.append(DocumentChunk(
                content="\n".join(slide_text),
                source_file=filename,
                file_type="pptx",
                slide_number=slide_num,
                section=slide_title,
                content_type="text"
            ))
    
    logger.info(f"  ✅ Extracted {len(chunks)} chunks from {filename}")
    
    return chunks


def load_text(file_path: str) -> list[DocumentChunk]:
    """Load and extract text from a plain text file."""
    filename = os.path.basename(file_path)
    
    logger.info(f"📄 Loading text file: {filename}")
    
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read().strip()
    
    # Clean noise
    content = clean_text(content)
    
    if content:
        logger.info(f"  ✅ Loaded {len(content)} characters from {filename}")
        return [DocumentChunk(
            content=content,
            source_file=filename,
            file_type="txt",
            content_type="text"
        )]
    return []


def load_document(file_path: str) -> list[DocumentChunk]:
    """
    Load a document based on its file extension.
    Returns a list of DocumentChunk objects with rich metadata.
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    
    loaders = {
        '.pdf': load_pdf,
        '.docx': load_docx,
        '.xlsx': load_excel,
        '.xls': load_excel,
        '.pptx': load_pptx,
        '.txt': load_text,
    }
    
    loader = loaders.get(ext)
    if loader is None:
        raise ValueError(f"Unsupported file type: {ext}")
    
    return loader(file_path)
